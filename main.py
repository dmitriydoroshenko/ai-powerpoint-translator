import os
import glob
from pptx import Presentation
import logging
from logger_config import setup_logging
from translator import translate_all
from file_utils import save_presentation

setup_logging()

def has_hlink(paragraph):
    """Проверяет наличие гиперссылок в параграфе с защитой от пустых rId."""
    for run in paragraph.runs:
        try:
            hlink = run.hyperlink
            if hlink is not None:
                if hasattr(hlink, 'rId') and hlink.rId:
                    return True
                if hlink.address is not None:
                    return True
        except Exception:
            return True
    return False

def extract_table_texts(shape):
    """Извлекает текст из объекта-таблицы."""
    texts = []
    locations = []
    
    if not shape.has_table:
        return texts, locations
        
    for row_idx, row in enumerate(shape.table.rows):
        for cell_idx, cell in enumerate(row.cells):
            if cell.text.strip():
                for para_idx, para in enumerate(cell.text_frame.paragraphs):
                    if para.text.strip() and not has_hlink(para):
                        texts.append(para.text.strip())
                        locations.append((row_idx, cell_idx, para_idx))
                    elif has_hlink(para):
                        logging.info(f"Пропуск параграфа в ячейке таблицы из-за гиперссылки: {para.text[:30]}...")
    
    return texts, locations

def collect_text_data(prs):
    """Извлекает весь текст и его местоположения из презентации."""
    all_texts = []
    text_locations = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Обработка обычных текстовых блоков
            if hasattr(shape, "text_frame") and shape.text.strip():
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    if para.text.strip():
                        if has_hlink(para):
                            logging.info(f"Пропуск параграфа с гиперссылкой: {para.text[:30]}...")
                            continue
                        
                        all_texts.append(para.text.strip())
                        text_locations.append(("paragraph", slide_idx, shape_idx, para_idx))
            
            # Обработка таблиц
            if shape.has_table:
                table_texts, table_locations = extract_table_texts(shape)
                all_texts.extend(table_texts)
                for (row_idx, cell_idx, para_idx) in table_locations:
                    text_locations.append(("table", slide_idx, shape_idx, row_idx, cell_idx, para_idx))
    
    return all_texts, text_locations

def apply_translations(prs, text_locations, translated_texts):
    """Обновляет презентацию переведенным текстом с сохранением форматирования."""
    for location, translated_text in zip(text_locations, translated_texts):
        if location[0] == "paragraph":
            _, slide_idx, shape_idx, para_idx = location
            shape = prs.slides[slide_idx].shapes[shape_idx]
            if hasattr(shape, "text_frame") and para_idx < len(shape.text_frame.paragraphs):
                paragraph = shape.text_frame.paragraphs[para_idx]
                
                # Сохранение исходного форматирования
                original_alignment = paragraph.alignment
                original_level = paragraph.level
                has_bullet = hasattr(paragraph, "format") and hasattr(paragraph.format, "bullet")

                orig_color = None
                if paragraph.runs:
                    try:
                        if hasattr(paragraph.runs[0].font.color, 'rgb'):
                            orig_color = paragraph.runs[0].font.color.rgb
                    except: pass

                original_font_sizes = [
                    run.font.size if hasattr(run, "font") and hasattr(run.font, "size") else None 
                    for run in paragraph.runs
                ]
                
                paragraph.text = translated_text
                
                for idx, run in enumerate(paragraph.runs):
                    run.font.name = "Microsoft YaHei"
                    if idx < len(original_font_sizes) and original_font_sizes[idx] is not None:
                        run.font.size = original_font_sizes[idx]
                    if orig_color:
                        run.font.color.rgb = orig_color
                
                paragraph.alignment = original_alignment
                paragraph.level = original_level
                if has_bullet and hasattr(paragraph, "format"):
                    try: paragraph.format.bullet.enable = True
                    except: logging.warning("Не удалось восстановить форматирование маркеров")

        elif location[0] == "table":
            _, slide_idx, shape_idx, row_idx, cell_idx, para_idx = location
            shape = prs.slides[slide_idx].shapes[shape_idx]
            if shape.has_table:
                cell = shape.table.rows[row_idx].cells[cell_idx]
                if hasattr(cell, "text_frame") and para_idx < len(cell.text_frame.paragraphs):
                    paragraph = cell.text_frame.paragraphs[para_idx]
                    
                    original_alignment = paragraph.alignment
                    original_level = paragraph.level
                    has_bullet = hasattr(paragraph, "format") and hasattr(paragraph.format, "bullet")
                    
                    orig_color = None
                    if paragraph.runs:
                        try:
                            if hasattr(paragraph.runs[0].font.color, 'rgb'):
                                orig_color = paragraph.runs[0].font.color.rgb
                        except: pass

                    original_font_sizes = [
                        run.font.size if hasattr(run, "font") and hasattr(run.font, "size") else None 
                        for run in paragraph.runs
                    ]
                    
                    paragraph.text = translated_text
                    
                    for idx, run in enumerate(paragraph.runs):
                        run.font.name = "Microsoft YaHei"
                        if idx < len(original_font_sizes) and original_font_sizes[idx] is not None:
                            run.font.size = original_font_sizes[idx]
                        if orig_color:
                            run.font.color.rgb = orig_color
                    
                    paragraph.alignment = original_alignment
                    paragraph.level = original_level
                    if has_bullet and hasattr(paragraph, "format"):
                        try: paragraph.format.bullet.enable = True
                        except: logging.warning("Не удалось восстановить форматирование маркеров в таблице")

def process_presentation(input_file):
    """Основной цикл обработки файла."""
    logging.info(f"Обработка файла: {input_file}")
    print(f"\nОбработки файла: {os.path.basename(input_file)}")
    
    try:
        prs = Presentation(input_file)
        
        # 1. Извлечение
        all_texts, text_locations = collect_text_data(prs)
        
        if not all_texts:
            logging.info(f"В файле {input_file} текст не найден")
            return
        
        # 2. Перевод
        translated_texts = translate_all(all_texts)
        
        # 3. Обновление
        apply_translations(prs, text_locations, translated_texts)
        
        # 4. Сохранение
        save_presentation(prs, input_file)
        
    except Exception as e:
        logging.error(f"Ошибка при обработке презентации {input_file}: {str(e)}")
        raise

def main():
    # Поиск всех файлов PPTX в директории input
    input_files = glob.glob('input/*.pptx')
    
    if not input_files:
        logging.warning("В директории input не найдено файлов PowerPoint")
        print("В директории input не найдено файлов PowerPoint")
        return
    
    for input_file in input_files:
        logging.info(f"\n=== Обработка файла: {input_file} ===")
        process_presentation(input_file)
        logging.info(f"Перевод файла {input_file} завершен")

if __name__ == "__main__":
    main()