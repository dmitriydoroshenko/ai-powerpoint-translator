import os
import glob
import json
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv
import logging
import time
from datetime import datetime
import httpx
import sys
import re

# Set up logging first, before any other imports
def setup_logging():
    # Create logs directory if it doesn't exist
    log_dir = 'SlideTranslateLog'
    os.makedirs(log_dir, exist_ok=True)
    
    # Create a timestamp for the log file
    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
    log_file = os.path.join(log_dir, f'{timestamp}.log')
    
    # Configure logging to only file
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler(log_file, encoding='utf-8')
        ]
    )
    return log_file

# Initialize logging first
log_file = setup_logging()
logging.info("Logging system initialized")

# Load environment variables
load_dotenv()


# Initialize OpenAI client
custom_http_client = httpx.Client() # Explicitly disable proxies if not needed

client = OpenAI(
    api_key=os.getenv('OPENAI_API_KEY'),
    http_client=custom_http_client # Truyền HTTP client tùy chỉnh vào
)

def batch_texts(texts, batch_size=30):
    """Group texts into batches for translation."""
    return [texts[i:i + batch_size] for i in range(0, len(texts), batch_size)]

def translate_batch(texts):
    """Translate a batch of texts from English to Simplified Chinese."""
    if not texts:
        return []
    
    payload = {f"item_{i}": text for i, text in enumerate(texts)}
    json_payload = json.dumps(payload, ensure_ascii=False)
    
    # Логируем запрос
    logging.info(f"=== Translation Request (Batch size: {len(texts)}) ===")
    for idx, text in enumerate(texts):
        logging.info(f"Input Item {idx}: {text}")
    
    try:
        SYSTEM_ROLE = (
            "You are a professional mobile game localizer (English to Simplified Chinese). "
            "Expertise: gaming terminology, UI/UX constraints, and mobile gaming slang. "
            "IMPORTANT: Preserve all special characters like vertical tabs (\\u000b), "
            "newlines (\\n), and specific spacing. Do not clean up the formatting. "
            "Task: Translate values to Simplified Chinese. Keep keys unchanged. "
            "Output: Return a valid JSON object."
        )

        response = client.chat.completions.create(
            model="gpt-5.2",
            messages=[
                {"role": "system", "content": SYSTEM_ROLE},
                {"role": "user", "content": f"Translate these items:\n{json_payload}"}
            ],
            response_format={ "type": "json_object" },
            temperature=0.3
        )
        
        raw_content = response.choices[0].message.content.strip()
        
        # Логируем сырой ответ от API
        logging.info("=== Raw API Response ===")
        logging.info(raw_content)
        
        translated_data = json.loads(raw_content)
        translations = []
        
        logging.info("=== Parsed Translations ===")
        for i in range(len(texts)):
            key = f"item_{i}"
            # Если ключа нет, логируем предупреждение и оставляем оригинал
            if key in translated_data:
                trans_text = translated_data[key].strip()
                translations.append(trans_text)
                logging.info(f"Item {i}: OK -> {trans_text}")
            else:
                logging.warning(f"Item {i}: MISSING in response! Using original text.")
                translations.append(texts[i])
        
        logging.info("=== Batch Processing Completed ===\n")

        # Логирование использования токенов:
        usage = response.usage
        logging.info(f"Tokens used - Prompt: {usage.prompt_tokens}, "
                    f"Completion: {usage.completion_tokens}, "
                    f"Total: {usage.total_tokens}")
        
        # Примерная стоимость для gpt-5.2 на момент 2026:
        # Цена за 1 млн токенов: $1.75 (input) / $14.00 (output)
        cost = (usage.prompt_tokens * 1.75 / 1_000_000) + (usage.completion_tokens * 14.00 / 1_000_000)
        logging.info(f"Estimated batch cost: ${cost:.4f}")

        return translations
        
    except Exception as e:
        logging.error(f"Error during translation: {str(e)}")
        print(f"\n[ERROR] Check logs: {str(e)}")
        raise

def save_presentation(prs, original_filename):
    """Save presentation with error handling and unique filename."""
    # Create output directory if it doesn't exist
    output_dir = 'output'
    os.makedirs(output_dir, exist_ok=True)
    
    # Generate base output filename
    base_name = os.path.basename(original_filename)
    name_without_ext = os.path.splitext(base_name)[0]
    
    # Try to save with different names if file exists or is locked
    counter = 1
    while True:
        if counter == 1:
            output_filename = os.path.join(output_dir, f"{name_without_ext}_cn.pptx")
        else:
            output_filename = os.path.join(output_dir, f"{name_without_ext}_cn_{counter}.pptx")
        
        try:
            prs.save(output_filename)
            logging.info(f"Successfully saved presentation to {output_filename}")
            return output_filename
        except PermissionError:
            logging.warning(f"Permission denied when saving to {output_filename}. File might be open in PowerPoint.")
            logging.info("Please close the file in PowerPoint if it's open.")
            counter += 1
            if counter > 5:  # Limit number of attempts
                raise Exception(f"Failed to save presentation after {counter-1} attempts. Please ensure the file is not open in PowerPoint.")
        except Exception as e:
            logging.error(f"Error saving presentation: {str(e)}")
            raise

def has_hlink(paragraph):
    """Проверяет наличие гиперссылок в параграфе с защитой от пустых rId."""
    for run in paragraph.runs:
        try:
            hlink = run.hyperlink
            if hlink is not None:
                if hasattr(hlink, 'rId') and hlink.rId:
                    return True
                if hlink.address is not None:
                    return True
        except Exception:
            return True
    return False

def extract_table_texts(shape):
    """Extract texts from a table shape."""
    texts = []
    locations = []
    
    if not shape.has_table:
        return texts, locations
        
    for row_idx, row in enumerate(shape.table.rows):
        for cell_idx, cell in enumerate(row.cells):
            if cell.text.strip():
                for para_idx, para in enumerate(cell.text_frame.paragraphs):
                    if para.text.strip() and not has_hlink(para):
                        texts.append(para.text.strip())
                        locations.append((row_idx, cell_idx, para_idx))
                    elif has_hlink(para):
                        logging.info(f"Skipping table cell paragraph with hyperlink: {para.text[:30]}...")
    
    return texts, locations

def split_text_by_paragraphs(text):
    """Split text into paragraphs, handling bullet points and line breaks."""
    # First, split by line breaks
    lines = text.split('\n')
    result = []
    current_text = ""
    
    for line in lines:
        line = line.strip()
        if not line:
            if current_text:
                result.append(current_text)
                current_text = ""
            continue
            
        # Check if line starts with bullet or numbering
        if re.match(r'^[•\-\*]|\d+[.)]', line):
            # This is likely a new bullet or numbered item
            if current_text:
                result.append(current_text)
            current_text = line
        elif current_text:
            # Check if previous line had a bullet and this is continuation
            if re.match(r'^[•\-\*]|\d+[.)]', current_text.split('\n')[0]):
                current_text += '\n' + line
            else:
                # Likely a separate paragraph
                result.append(current_text)
                current_text = line
        else:
            current_text = line
    
    # Add the last paragraph if any
    if current_text:
        result.append(current_text)
        
    return result

def process_presentation(input_file):
    """Process a PowerPoint presentation, translating text from English to Simplified Chinese."""
    logging.info(f"Processing {input_file}")
    
    try:
        prs = Presentation(input_file)
        all_texts = []
        text_locations = []
        
        # Extract all texts that need translation
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Handle regular text shapes
                if hasattr(shape, "text_frame") and shape.text.strip():
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        if para.text.strip():
                            if has_hlink(para):
                                logging.info(f"Skipping paragraph with hyperlink: {para.text[:30]}...")
                                continue
                                
                            all_texts.append(para.text.strip())
                            text_locations.append(("paragraph", slide_idx, shape_idx, para_idx))
                
                # Handle tables
                if shape.has_table:
                    table_texts, table_locations = extract_table_texts(shape)
                    all_texts.extend(table_texts)
                    for (row_idx, cell_idx, para_idx) in table_locations:
                        text_locations.append(("table", slide_idx, shape_idx, row_idx, cell_idx, para_idx))
        
        if not all_texts:
            logging.info(f"No text found in {input_file}")
            return
        
        # Translate texts in batches
        translated_texts = []
        batches = batch_texts(all_texts)
        
        print(f"\nTranslating {os.path.basename(input_file)}:")
        for i, batch in enumerate(batches):
            progress = (i + 1) / len(batches) * 100
            sys.stdout.write(f"\rProgress: [{int(progress)}%] Batch {i+1}/{len(batches)}")
            sys.stdout.flush()
            
            logging.info(f"Translating batch {i+1}/{len(batches)} (size: {len(batch)} texts)")
            translations = translate_batch(batch)
            translated_texts.extend(translations)
            
            if i < len(batches) - 1:
                time.sleep(2)
        print("\nTranslation completed!")
        
        # Update presentation with translations
        for location, translated_text in zip(text_locations, translated_texts):
            if location[0] == "paragraph":
                _, slide_idx, shape_idx, para_idx = location
                shape = prs.slides[slide_idx].shapes[shape_idx]
                if hasattr(shape, "text_frame") and para_idx < len(shape.text_frame.paragraphs):
                    paragraph = shape.text_frame.paragraphs[para_idx]
                    
                    # Store original formatting
                    original_alignment = paragraph.alignment
                    original_level = paragraph.level
                    has_bullet = False
                    if hasattr(paragraph, "format") and hasattr(paragraph.format, "bullet"):
                        has_bullet = True

                    # Extract original color before clearing text
                    orig_color = None
                    if paragraph.runs:
                        try:
                            if hasattr(paragraph.runs[0].font.color, 'rgb'):
                                orig_color = paragraph.runs[0].font.color.rgb
                        except:
                            pass

                    # Store original font sizes before updating text
                    original_font_sizes = []
                    for run in paragraph.runs:
                        if hasattr(run, "font") and hasattr(run.font, "size"):
                            original_font_sizes.append(run.font.size)
                        else:
                            original_font_sizes.append(None)  # None means use default
                    
                    paragraph.text = translated_text
                    
                    # Set font to Microsoft YaHei for all runs in the paragraph while keeping original size
                    for idx, run in enumerate(paragraph.runs):
                        run.font.name = "Microsoft YaHei"
                        # If we have stored a font size and have enough runs, use the original
                        if idx < len(original_font_sizes) and original_font_sizes[idx] is not None:
                            run.font.size = original_font_sizes[idx]
                        
                        # Apply original color if found
                        if orig_color:
                            run.font.color.rgb = orig_color
                    
                    # Restore original formatting
                    paragraph.alignment = original_alignment
                    paragraph.level = original_level
                    if has_bullet and hasattr(paragraph, "format"):
                        try:
                            paragraph.format.bullet.enable = True
                        except:
                            # If bullet restoration fails, log but continue
                            logging.warning("Failed to restore bullet formatting")
            
            elif location[0] == "text":
                _, slide_idx, shape_idx, _ = location
                shape = prs.slides[slide_idx].shapes[shape_idx]
                shape.text = translated_text
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "Microsoft YaHei"
                
            elif location[0] == "table":
                _, slide_idx, shape_idx, row_idx, cell_idx, para_idx = location
                shape = prs.slides[slide_idx].shapes[shape_idx]
                if shape.has_table:
                    cell = shape.table.rows[row_idx].cells[cell_idx]
                    if hasattr(cell, "text_frame") and para_idx < len(cell.text_frame.paragraphs):
                        paragraph = cell.text_frame.paragraphs[para_idx]
                        
                        # Store original formatting
                        original_alignment = paragraph.alignment
                        original_level = paragraph.level
                        has_bullet = False
                        if hasattr(paragraph, "format") and hasattr(paragraph.format, "bullet"):
                            has_bullet = True
                        
                        # Extract original color for tables
                        orig_color = None
                        if paragraph.runs:
                            try:
                                if hasattr(paragraph.runs[0].font.color, 'rgb'):
                                    orig_color = paragraph.runs[0].font.color.rgb
                            except:
                                pass

                        # Store original font sizes before updating text
                        original_font_sizes = []
                        for run in paragraph.runs:
                            if hasattr(run, "font") and hasattr(run.font, "size"):
                                original_font_sizes.append(run.font.size)
                            else:
                                original_font_sizes.append(None)  # None means use default
                        
                        paragraph.text = translated_text
                        
                        # Set font to Microsoft YaHei for all runs in the paragraph while keeping original size
                        for idx, run in enumerate(paragraph.runs):
                            run.font.name = "Microsoft YaHei"
                            # If we have stored a font size and have enough runs, use the original
                            if idx < len(original_font_sizes) and original_font_sizes[idx] is not None:
                                run.font.size = original_font_sizes[idx]
                            
                            # Apply color to table text
                            if orig_color:
                                run.font.color.rgb = orig_color
                        
                        # Restore original formatting
                        paragraph.alignment = original_alignment
                        paragraph.level = original_level
                        if has_bullet and hasattr(paragraph, "format"):
                            try:
                                paragraph.format.bullet.enable = True
                            except:
                                # If bullet restoration fails, log but continue
                                logging.warning("Failed to restore bullet formatting")
        
        # Save translated presentation with error handling
        save_presentation(prs, input_file)
        
    except Exception as e:
        logging.error(f"Error processing presentation {input_file}: {str(e)}")
        raise

def main():
    # Find all PPTX files in the input directory
    input_files = glob.glob('input/*.pptx')
    
    if not input_files:
        logging.warning("No PowerPoint files found in the input directory")
        return
    
    for input_file in input_files:
        logging.info(f"\n=== Processing file: {input_file} ===")
        process_presentation(input_file)
        logging.info(f"Completed translation of {input_file}")

if __name__ == "__main__":
    main()